import os
import pdfplumber
from docx import Document
from pptx import Presentation
from pdf2image import convert_from_path
import pytesseract
import fitz  # PyMuPDF
from PIL import Image


# ============================================================
# 📘 Master extractor — detects file type
# ============================================================
def extract_text_and_images(file_path: str, output_dir: str = "extracted_images"):
    """
    Extract text + image file paths from PDF, DOCX, or PPTX.
    Returns: {"text": str, "images": [list of image paths]}
    """
    ext = os.path.splitext(file_path)[1].lower()
    text, images = "", []

    if ext == ".pdf":
        text = _parse_pdf(file_path)
        images = _extract_images_from_pdf(file_path, output_dir)
    elif ext == ".docx":
        text = _parse_docx(file_path)
        images = _extract_images_from_docx(file_path, output_dir)
    elif ext == ".pptx":
        text = _parse_pptx(file_path)
        images = _extract_images_from_pptx(file_path, output_dir)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    text = " ".join(text.split())  # clean whitespace
    return {"text": text, "images": images}


# ------------------------------------------------------------
# 🔹 PDF text (with OCR fallback)
# ------------------------------------------------------------
def _parse_pdf(file_path: str) -> str:
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"

    # OCR fallback if no selectable text
    if not text.strip():
        print("🧩 No text detected — performing OCR...")
        images = convert_from_path(file_path)
        ocr_texts = [pytesseract.image_to_string(img) for img in images]
        text = "\n".join(ocr_texts)
    return text


# ------------------------------------------------------------
# 🔹 PDF image extraction
# ------------------------------------------------------------
def _extract_images_from_pdf(file_path: str, output_dir: str):
    os.makedirs(output_dir, exist_ok=True)
    doc = fitz.open(file_path)
    paths = []
    for page_index, page in enumerate(doc):
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            image_path = os.path.join(output_dir, f"page{page_index+1}_img{img_index+1}.{image_ext}")
            with open(image_path, "wb") as f:
                f.write(image_bytes)
            paths.append(image_path)
    return paths


# ------------------------------------------------------------
# 🔹 DOCX text & image extraction
# ------------------------------------------------------------
def _parse_docx(file_path: str) -> str:
    doc = Document(file_path)
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paras)


def _extract_images_from_docx(file_path: str, output_dir: str):
    from zipfile import ZipFile
    os.makedirs(output_dir, exist_ok=True)
    image_paths = []
    with ZipFile(file_path, "r") as z:
        for f in z.namelist():
            if f.startswith("word/media/") and (f.endswith(".png") or f.endswith(".jpg") or f.endswith(".jpeg")):
                img_name = os.path.basename(f)
                dest = os.path.join(output_dir, img_name)
                with open(dest, "wb") as out:
                    out.write(z.read(f))
                image_paths.append(dest)
    return image_paths


# ------------------------------------------------------------
# 🔹 PPTX text & image extraction
# ------------------------------------------------------------
def _parse_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)


def _extract_images_from_pptx(file_path: str, output_dir: str):
    from zipfile import ZipFile
    os.makedirs(output_dir, exist_ok=True)
    image_paths = []
    with ZipFile(file_path, "r") as z:
        for f in z.namelist():
            if f.startswith("ppt/media/") and (f.endswith(".png") or f.endswith(".jpg") or f.endswith(".jpeg")):
                img_name = os.path.basename(f)
                dest = os.path.join(output_dir, img_name)
                with open(dest, "wb") as out:
                    out.write(z.read(f))
                image_paths.append(dest)
    return image_paths
